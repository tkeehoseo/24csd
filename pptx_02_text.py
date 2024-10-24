from pptx import Presentation
from pptx.util import Cm

# 새 프레젠테이션 생성
prs = Presentation()

# 슬라이드 추가 (빈 슬라이드 레이아웃 사용)
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

# 텍스트 상자 추가
left = top = Cm(2)
width = height = Cm(4)
txBox = slide.shapes.add_textbox(left, top, width, height)

# 텍스트 상자 내부에 텍스트 추가
tf = txBox.text_frame
tf.text = "Python-pptx를 사용하여 텍스트 상자 추가하기"

# 추가 텍스트 줄 추가
p = tf.add_paragraph()
p.text = "또 다른 문장을 추가합니다."
p.font.bold = True  # 텍스트를 굵게

prs.save('text_box_example.pptx')
