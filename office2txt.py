import sys
import os
from pptx import Presentation
from docx import Document
import openpyxl

def extract_text_from_pptx(input_pptx, output_txt):
    """PPTX 파일에서 텍스트를 추출하여 TXT 파일로 저장"""
    prs = Presentation(input_pptx)
    with open(output_txt, 'w', encoding='utf-8') as f:
        for i, slide in enumerate(prs.slides):
            f.write(f"--- Slide {i + 1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    f.write(shape.text.strip() + "\n")
            f.write("\n")

def extract_text_from_docx(input_docx, output_txt):
    """DOCX 파일에서 텍스트를 추출하여 TXT 파일로 저장"""
    doc = Document(input_docx)
    with open(output_txt, 'w', encoding='utf-8') as f:
        for para in doc.paragraphs:
            f.write(para.text.strip() + "\n")
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                f.write("\t".join(row_text) + "\n")

def extract_text_from_xlsx(input_xlsx, output_txt):
    """XLSX 파일에서 텍스트를 추출하여 TXT 파일로 저장"""
    wb = openpyxl.load_workbook(input_xlsx)
    with open(output_txt, 'w', encoding='utf-8') as f:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            f.write(f"--- Worksheet: {sheet_name} ---\n")
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else '' for cell in row]
                f.write("\t".join(row_data) + "\n")
            f.write("\n")

def main():
    if len(sys.argv) != 3:
        print("사용법: python office2txt.py <input_file> <output_txt>")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_txt = sys.argv[2]
    
    # 입력 파일의 확장자 확인
    _, file_extension = os.path.splitext(input_file)
    
    # 확장자에 따른 처리
    if file_extension.lower() == '.pptx':
        extract_text_from_pptx(input_file, output_txt)
        print(f"PPTX 파일의 텍스트가 {output_txt} 파일로 저장되었습니다.")
    elif file_extension.lower() == '.docx':
        extract_text_from_docx(input_file, output_txt)
        print(f"DOCX 파일의 텍스트가 {output_txt} 파일로 저장되었습니다.")
    elif file_extension.lower() == '.xlsx':
        extract_text_from_xlsx(input_file, output_txt)
        print(f"XLSX 파일의 텍스트가 {output_txt} 파일로 저장되었습니다.")
    else:
        print("지원하지 않는 파일 형식입니다. .pptx, .docx, .xlsx 파일만 지원합니다.")
        sys.exit(1)

if __name__ == "__main__":
    main()
