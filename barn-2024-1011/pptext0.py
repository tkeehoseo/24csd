from pptx import Presentation
prs = Presentation("./ppt/sample.pptx")
text_runs = []
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
for s in text_runs:
    print(s, end=' ')