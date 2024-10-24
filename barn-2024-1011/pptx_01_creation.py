from pptx import Presentation

# 새 프레젠테이션 생성
prs = Presentation()

# 슬라이드 추가 (제목 슬라이드 레이아웃 사용)
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)

# 제목과 부제목 설정
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, PowerPoint!"
subtitle.text = "python-pptx 라이브러리 사용 예제"

# 프레젠테이션 저장
prs.save('example.pptx')
