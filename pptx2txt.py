import sys
from pptx import Presentation

def extract_text_from_pptx(input_pptx, output_txt):
    # PowerPoint 파일 열기
    prs = Presentation(input_pptx)
    
    # 출력할 텍스트 파일 열기
    with open(output_txt, 'w', encoding='utf-8') as f:
        # 슬라이드별로 텍스트 추출
        for i, slide in enumerate(prs.slides):
            f.write(f"--- Slide {i + 1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    f.write(shape.text.strip() + "\n")
            f.write("\n")  # 슬라이드 구분을 위한 빈 줄

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("사용법: python ppt2txt.py <input_pptx> <output_txt>")
    else:
        input_pptx = sys.argv[1]
        output_txt = sys.argv[2]
        extract_text_from_pptx(input_pptx, output_txt)
        print(f"텍스트가 {output_txt} 파일에 저장되었습니다.")
